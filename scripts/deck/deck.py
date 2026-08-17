"""Saniti 디자인 시스템을 python-pptx로 구현한다. 좌표는 1920x1080 기준틀의 px."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageFont, ImageChops
from pathlib import Path
import copy, os, functools

# ── 텍스트 계측 ──────────────────────────────────────────────────────────
_F = os.path.expanduser("~/.fonts")
_FONTS = {
    (SANS_ := "Inter", False): f"{_F}/Inter-Regular.ttf",
    ("Inter", True): f"{_F}/Inter-Bold.ttf",
    ("IBM Plex Mono", False): f"{_F}/IBMPlexMono-Regular.ttf",
    ("IBM Plex Mono", True): f"{_F}/IBMPlexMono-SemiBold.ttf",
    # 한글 — 측정용 실물. 없으면 줄바꿈 계산이 통째로 틀어진다.
    ("NanumGothic", False): f"{_F}/NanumGothic.ttf",
    ("NanumGothic", True): f"{_F}/NanumGothicBold.ttf",
}

@functools.lru_cache(maxsize=512)
def _font(name, bold, px):
    path = _FONTS[(name, bold)]
    if not os.path.exists(path):
        # 조용한 실패 금지 — 폰트가 없으면 폭 계산이 엉켜 글자가 카드를 넘친다.
        raise FileNotFoundError(
            f"측정용 폰트가 없다: {path}\n"
            "  한글 덱: python -c \"import koreanize_matplotlib,glob,os,shutil;"
            "[shutil.copy(f,os.path.expanduser('~/.fonts/')) for f in "
            "glob.glob(os.path.dirname(koreanize_matplotlib.__file__)+'/**/*.ttf',recursive=True)]\"")
    return ImageFont.truetype(path, int(round(px)))

def text_w(txt, size_pt, font=None, bold=False, spacing=0.0):
    """1920 기준틀 px 폭. 1pt = 2px(144dpi 기준틀). font=None이면 현재 테마의 SANS."""
    f = _font(font or SANS, bold, size_pt * 2)
    return f.getlength(txt) + spacing * 2 * len(txt)

def wrap_lines(txt, w_px, size_pt, font=None, bold=False, first_indent=0.0):
    """주어진 폭에서 몇 줄이 되는지. 실제 폰트 메트릭으로 잰다."""
    words, lines, cur = txt.split(), 0, ""
    avail = w_px - first_indent
    for wd in words:
        trial = (cur + " " + wd).strip()
        if text_w(trial, size_pt, font, bold) <= avail or not cur:
            cur = trial
        else:
            lines += 1; cur = wd; avail = w_px
    return lines + (1 if cur else 0)

LINE_EM = 1.212          # Inter의 단일 행간(em) — 렌더러가 쓰는 값
SAFE    = 8              # 줄바꿈 계산에서 빼두는 여유 폭(px)

def line_px(size_pt, line=1.35):
    return size_pt * 2 * LINE_EM * line

def block_h(txt, w_px, size_pt, line=1.35):
    """본문 한 덩어리의 높이(px)."""
    n = wrap_lines(txt, w_px - SAFE, size_pt) if isinstance(txt, str) else len(txt)
    return n * line_px(size_pt, line)

EMU_PX = 6350                      # 1px(1920 기준틀) = 1/144 inch
def E(px): return Emu(int(round(px * EMU_PX)))

# ── 테마 ────────────────────────────────────────────────────────────────
# 환경변수 DECK_THEME=light 면 **흰 바탕·검정 글씨**. 기본은 기존 다크(옛 덱 재현용).
# 색 상수는 **아래 함수들의 기본 인자로 캡처**되므로 반드시 클래스 정의 **전에** 정해야 한다.
#   이름은 그대로 둔다(WHITE/ASH/…) — 200군데 호출부를 안 건드리려는 것이다.
#   라이트에서 WHITE는 "가장 진한 글씨색", CANVAS는 "바탕"이라는 **역할 이름**으로 읽는다.
_LIGHT = os.environ.get("DECK_THEME", "dark").lower() == "light"

if _LIGHT:
    CANVAS = RGBColor(0xFF, 0xFF, 0xFF)   # 바탕 — 흰색
    SOFT   = RGBColor(0xFA, 0xFA, 0xFA)   # 카드 바탕 — 거의 흰색
    WHITE  = RGBColor(0x0A, 0x0A, 0x0A)   # 가장 진한 글씨 — 검정
    ASH    = RGBColor(0x2E, 0x2E, 0x2E)   # 본문
    MUTE   = RGBColor(0x6B, 0x6B, 0x6B)   # 보조
    HAIR   = RGBColor(0xD8, 0xD8, 0xD8)   # 실선
    BRAND  = RGBColor(0xB3, 0x4A, 0x28)   # 강조(드물게)
    LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)
else:
    CANVAS = RGBColor(0x0B, 0x0B, 0x0B); SOFT = RGBColor(0x1C, 0x1C, 0x1C)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF); ASH  = RGBColor(0xB9, 0xB9, 0xB9)
    MUTE   = RGBColor(0x79, 0x79, 0x79); HAIR = RGBColor(0x35, 0x35, 0x35)
    BRAND  = RGBColor(0xD9, 0x77, 0x57); LIGHT = RGBColor(0xFF, 0xFF, 0xFF)

# 한글 덱은 나눔고딕. 뷰어(윈도우·맥)에 없을 수 있으므로 PPTX에는 이 이름이 박히고,
# 없으면 뷰어가 대체 폰트를 쓴다. **측정은 여기 있는 실제 파일로 한다.**
if os.environ.get("DECK_LANG", "en").lower() == "ko":
    SANS = "NanumGothic"; MONO = "NanumGothic"
else:
    SANS = "Inter"; MONO = "IBM Plex Mono"

# 고정 앵커
ML, MR = 100, 1820; CW = MR - ML
Y_LABEL, Y_TITLE, Y_SUB, Y_RULE = 62, 96, 180, 240
BODY_TOP, BODY_BOT = 276, 985
Y_FOOT = 1030
GAP = 24; PAD = 28


class Deck:
    def __init__(self, name):
        self.p = Presentation()
        self.p.slide_width, self.p.slide_height = E(1920), E(1080)
        self.blank = self.p.slide_layouts[6]
        self.name = name
        self.n = 0

    # ── 원시 요소 ────────────────────────────────────────────────────────
    def _tb(self, s, x, y, w, h, text, size, color, font=SANS, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0.0,
            caps=False, line=None, wrap=True):
        box = s.shapes.add_textbox(E(x), E(y), E(w), E(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        lines = text if isinstance(text, list) else [text]
        for i, ln in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            if line:
                para.line_spacing = line
            r = para.add_run()
            r.text = ln.upper() if caps else ln
            f = r.font
            f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
            if spacing:
                r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
            for rr in (r,):
                rPr = rr._r.get_or_add_rPr()
                for tag in ("a:latin", "a:cs", "a:ea"):
                    e = rPr.find(qn(tag))
                    if e is None:
                        e = rPr.makeelement(qn(tag), {}); rPr.append(e)
                    e.set("typeface", font)
        return box

    def rect(self, s, x, y, w, h, fill=SOFT, line=HAIR, lw=1.0, radius=12):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            E(x), E(y), E(w), E(h))
        if radius:
            shp.adjustments[0] = min(0.5, radius / min(w, h))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line; shp.line.width = Pt(lw)
        shp.shadow.inherit = False
        shp.text_frame.text = ""
        return shp

    def hline(self, s, x, y, w, color=HAIR, lw=1.0):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(1.2))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background(); shp.shadow.inherit = False
        return shp

    # ── 슬라이드 ─────────────────────────────────────────────────────────
    def slide(self, chapter, title, subtitle, light=False):
        self.n += 1
        s = self.p.slides.add_slide(self.blank)
        bg = s.background.fill; bg.solid()
        bg.fore_color.rgb = LIGHT if light else CANVAS
        ink = RGBColor(0x0B, 0x0B, 0x0B) if light else WHITE
        sub = RGBColor(0x55, 0x55, 0x55) if light else ASH
        mut = RGBColor(0x77, 0x77, 0x77) if light else MUTE
        self._tb(s, ML, Y_LABEL, CW, 22, chapter, 9, mut, caps=True, spacing=1.4)
        if wrap_lines(title, CW - 4, 27) > 1:
            raise ValueError(f"title wraps to 2 lines (slide {self.n}): "
                             f"{text_w(title, 27):.0f}px > {CW}px\n  {title}")
        self._tb(s, ML, Y_TITLE, CW, 70, title, 27, ink, spacing=-0.5)
        self._tb(s, ML, Y_SUB, CW, 40, subtitle, 12.5, sub, spacing=-0.2)
        self.hline(s, ML, Y_RULE, CW, HAIR if not light else RGBColor(0xDD, 0xDD, 0xDD))
        self._tb(s, ML, Y_FOOT, CW * 0.6, 20, f"{self.name}  ·  {chapter}", 8.5, mut, spacing=1.0)
        self._tb(s, ML + CW * 0.6, Y_FOOT, CW * 0.4, 20, f"{self.n:02d}", 8.5, mut,
                 align=PP_ALIGN.RIGHT, spacing=1.0)
        s._page = self.n
        return s

    # ── 카드 구성요소 ────────────────────────────────────────────────────
    def card(self, s, x, y, w, h, accent=False, fill=SOFT):
        # 강조 테두리는 쓰지 않는다 — 카드 경계는 전부 hairline.
        return self.rect(s, x, y, w, h, fill=fill, line=HAIR)

    def eyebrow(self, s, x, y, text, w=800, color=MUTE):
        if text_w(text.upper(), 8.5, bold=True, spacing=1.1) > w:
            raise ValueError(f"eyebrow wraps (slide {self.n}): "
                             f"{text_w(text.upper(), 8.5, bold=True, spacing=1.1):.0f} > {w:.0f}px\n  {text}")
        return self._tb(s, x, y, w, 16, text, 8.5, color, bold=True, caps=True, spacing=1.1)

    def cardtitle(self, s, x, y, text, w=800, color=WHITE, size=13):
        return self._tb(s, x, y, w, 24, text, size, color, spacing=-0.2)

    def body(self, s, x, y, w, text, color=ASH, size=10.5, h=None, line=1.35,
             fit_h=None, tag=""):
        """fit_h를 주면 그 안에 들어가는지 확인하고, 안 되면 크기를 줄이거나 예외."""
        if fit_h is not None:
            for sz in (size, 10.2, 10.0):
                if block_h(text, w, sz, line) + 6 <= fit_h:
                    size = sz; break
            else:
                raise ValueError(
                    f"body overflows [{tag}]: need "
                    f"{block_h(text, w, 10.0, line):.0f}px, have {fit_h:.0f}px\n  {str(text)[:90]}")
        hh = h or block_h(text, w, size, line) + 6
        self._tb(s, x, y, w, hh, text, size, color, line=line)
        return hh

    def caption(self, s, x, bottom, w, text, color=MUTE):
        """카드 아래 패딩선(bottom)에 바닥을 맞춘다. 반환 = 캡션 윗변 y."""
        h = block_h(text, w, 9.5, 1.25) + 4
        self._tb(s, x, bottom - h, w, h, text, 9.5, color,
                 anchor=MSO_ANCHOR.BOTTOM, line=1.25)
        return bottom - h

    def symbols(self, s, x, y, w, items, size=9.5):
        """수식 아래 기호 정의 — 한 줄에 여러 개."""
        txt = "   ".join(f"{k}  {v}" for k, v in items)
        return self._tb(s, x, y, w, 20, txt, size, MUTE)

    def equation(self, s, path, x, y, h=72):
        """수식은 높이로 맞춘다(55~90px). 카드 좌측 패딩선에 붙인다."""
        return self.image(s, path, x, y, h=h)

    def deflist(self, s, x, y, w, items, size=10.5, pitch=46, color=ASH,
                lead_color=WHITE, fit_h=None, tag=""):
        if fit_h is not None:
            for pt in (pitch, 44, 42, 40, 38, 36):
                if self.deflist_h(w, items, size, pt) <= fit_h:
                    pitch = pt; break
            else:
                for sz in (10.2, 10.0):
                    if self.deflist_h(w, items, sz, 36) <= fit_h:
                        size, pitch = sz, 36; break
                else:
                    raise ValueError(
                        f"deflist overflows [{tag}]: need "
                        f"{self.deflist_h(w, items, 10.0, 36):.0f}px, have {fit_h:.0f}px "
                        f"({len(items)} items)")
        """— **Term**: explanation  — 걸개 들여쓰기를 pPr(marL/indent)로 준다."""
        box = s.shapes.add_textbox(E(x), E(y), E(w), E(len(items) * pitch + 10))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, (term, expl) in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.line_spacing = 1.34
            para.space_after = Pt(max(0.0, (pitch - line_px(size, 1.34))) / 2.0)
            # 걸개: "—  <Term>: " 폭만큼 marL, 그만큼 음수 indent
            lead = f"—  {term}: "
            widthpx = len(lead) * size * 0.52
            pPr = para._p.get_or_add_pPr()
            pPr.set("marL", str(int(widthpx * 6350)))
            pPr.set("indent", str(-int(widthpx * 6350)))
            for txt, col, bold in (("—  ", MUTE, False), (term + ":", lead_color, True),
                                   (" " + expl, color, False)):
                r = para.add_run(); r.text = txt
                r.font.name = SANS; r.font.size = Pt(size)
                r.font.bold = bold; r.font.color.rgb = col
                rPr = r._r.get_or_add_rPr()
                for tag in ("a:latin", "a:cs", "a:ea"):
                    e = rPr.find(qn(tag))
                    if e is None:
                        e = rPr.makeelement(qn(tag), {}); rPr.append(e)
                    e.set("typeface", SANS)
        return box

    @staticmethod
    def deflist_h(w, items, size=10.5, pitch=46):
        """정의 목록이 실제로 차지할 높이."""
        lh = line_px(size, 1.34)
        gap = max(0.0, pitch - lh)
        tot = 0.0
        for term, expl in items:
            ind = (text_w("—  ", size) + text_w(term + ":", size, bold=True)
                   + text_w(" ", size))
            avail = w - ind - SAFE
            words, n, cur = expl.split(), 1, ""
            for wd in words:
                t = (cur + " " + wd).strip()
                if text_w(t, size) <= avail or not cur:
                    cur = t
                else:
                    n += 1; cur = wd
            tot += n * lh + gap
        return tot - gap

    def panel(self, s, x, y, w, h, eyebrow=None, caption=None, accent=False,
              fill=SOFT, light=False):
        """카드 + 눈썹 + 캡션 밴드를 한 번에. 반환 = 쓸 수 있는 본문 영역 (x,y,w,h)."""
        self.card(s, x, y, w, h, accent=accent, fill=fill)
        cx, cw = x + PAD, w - 2 * PAD
        cy = y + PAD
        if eyebrow:
            self.eyebrow(s, cx, cy, eyebrow, w=cw,
                         color=RGBColor(0x77,0x77,0x77) if light else MUTE)
            cy += 30
        bottom = y + h - PAD
        if caption:
            bottom = self.caption(s, cx, bottom, cw, caption,
                                  color=RGBColor(0x77,0x77,0x77) if light else MUTE) - 12
        return (cx, cy, cw, bottom - cy)

    def stat(self, s, x, y, w, value, label, color=WHITE, size=44, fit_h=None,
             tag=""):
        vh = size * 2 * 1.15
        lh = block_h(label.upper(), w, 8.5, 1.3)
        if fit_h is not None and vh + lh + 10 > fit_h:
            raise ValueError(f"stat overflows [{tag}]: need {vh+lh+10:.0f}px, "
                             f"have {fit_h:.0f}px — shorten “{label}”")
        self._tb(s, x, y, w, vh, value, size, color, spacing=-1.0)
        self._tb(s, x, y + vh + 4, w, lh + 6, label, 8.5, MUTE, bold=True,
                 caps=True, spacing=1.1, line=1.3)
        return vh + lh + 10

    # ── 이미지 ───────────────────────────────────────────────────────────
    def image(self, s, path, x, y, w=None, h=None, align="left"):
        iw, ih = Image.open(path).size
        if w and not h: h = w * ih / iw
        elif h and not w: w = h * iw / ih
        elif not w and not h: raise ValueError("w or h required")
        if align == "center": x = x - w / 2
        s.shapes.add_picture(str(path), E(x), E(y), E(w), E(h))
        return (x, y, w, h)

    def fit(self, s, path, x, y, w, h, align="center", pad=0):
        """박스 안에 비율 유지로 맞춘다. 반환 = 실제 배치 사각형."""
        iw, ih = Image.open(path).size
        bw, bh = w - 2 * pad, h - 2 * pad
        sc = min(bw / iw, bh / ih)
        dw, dh = iw * sc, ih * sc
        dx = x + pad + (0 if align == "left" else (bw - dw) / 2)
        dy = y + pad + (bh - dh) / 2
        s.shapes.add_picture(str(path), E(dx), E(dy), E(dw), E(dh))
        return (dx, dy, dw, dh)

    @staticmethod
    def trim(path):
        """결과 PNG의 흰 여백을 잘라 캐시한다 — 카드 안에서 그림이 작아 보이지 않게."""
        src = Path(str(path))
        dst = Path("assets/trim") / (src.parent.parent.name + "__" + src.name)
        if not dst.exists():
            im = Image.open(src).convert("RGB")
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bbox = ImageChops.difference(im, bg).getbbox()
            if bbox:
                pad = 6
                bbox = (max(0, bbox[0] - pad), max(0, bbox[1] - pad),
                        min(im.width, bbox[2] + pad), min(im.height, bbox[3] + pad))
                im = im.crop(bbox)
            dst.parent.mkdir(parents=True, exist_ok=True)
            im.save(dst)
        return dst

    def plot(self, s, path, x, y, w, h, pad=12, align="center"):
        """흰 카드를 그림 비율에 맞춰 만든다 — 카드 안에 빈 흰 여백이 남지 않게."""
        src = self.trim(path)
        iw, ih = Image.open(src).size
        sc = min((w - 2 * pad) / iw, (h - 2 * pad) / ih)
        dw, dh = iw * sc, ih * sc
        cw_, ch_ = dw + 2 * pad, dh + 2 * pad
        cx = x if align == "left" else x + (w - cw_) / 2
        cy = y + (h - ch_) / 2
        self.rect(s, cx, cy, cw_, ch_, fill=LIGHT, line=HAIR, radius=12)
        s.shapes.add_picture(str(src), E(cx + pad), E(cy + pad), E(dw), E(dh))
        return (cx, cy, cw_, ch_)

    def save(self, path):
        self.p.save(path)
        return path
